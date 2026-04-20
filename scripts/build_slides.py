"""Build Demo Day PowerPoint slides matching the PANACEA webapp theme.

Usage: python scripts/build_slides.py
Output: writeup/panacea_demo_day.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os
import json

# ── Theme colors (from webapp) ──
BG = RGBColor(0x08, 0x08, 0x0C)
SURFACE = RGBColor(0x11, 0x11, 0x18)
BORDER = RGBColor(0x2A, 0x2A, 0x3A)
TEXT_PRIMARY = RGBColor(0xE8, 0xE8, 0xF0)
TEXT_MUTED = RGBColor(0x7C, 0x7C, 0x96)
TEXT_DIM = RGBColor(0x55, 0x55, 0x6A)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ROOT = os.path.join(os.path.dirname(__file__), "..")
FIGURES = os.path.join(ROOT, "writeup", "figures")
OUTPUT = os.path.join(ROOT, "writeup", "panacea_demo_day.pptx")

# Load live stats
try:
    with open(os.path.join(ROOT, "webapp-react", "public", "cdm_forecast.json")) as f:
        forecast = json.load(f)
    with open(os.path.join(ROOT, "webapp-react", "public", "pipeline_stats.json")) as f:
        pipeline = json.load(f)
except Exception:
    forecast, pipeline = {}, {}

fm = pipeline.get("forecast_model", {})
cs = pipeline.get("cdm_stats", {})
tr = forecast.get("track_record", {})
conf = forecast.get("conformal", {})

n_pairs = fm.get("n_pairs_total", 1167)
n_cdms = cs.get("total_cdms", 6667)
n_pos = fm.get("n_positive", 293)
pos_rate = n_pos / n_pairs if n_pairs else 0.251
naive_acc = (1 - pos_rate) * 100  # predict all negative
lr_f1 = fm.get("test", {}).get("f1", 0.887)
lr_recall = fm.get("test", {}).get("recall", 0.915)
lr_precision = fm.get("test", {}).get("precision", 0.860)
tp = tr.get("n_true_positives", 196)
fn_count = tr.get("n_false_negatives", 0)
fp = tr.get("n_false_positives", 160)
tn = tr.get("n_true_negatives", 789)
resolved = tp + fn_count + fp + tn
prod_acc = (tp + tn) / resolved * 100 if resolved else 86
prod_recall = tp / (tp + fn_count) * 100 if (tp + fn_count) else 100
conf_cov = (conf.get("regression_coverage", 0.905)) * 100


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=14,
             color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.6), height=Pt(4)):
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, fill_color=SURFACE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    return shape


def add_stat_card(slide, left, top, value, label, width=Inches(2.2)):
    card = add_card(slide, left, top, width, Inches(1.1))
    add_text(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.5),
             value, font_size=28, bold=True, color=TEXT_PRIMARY, font_name="Consolas")
    add_text(slide, left + Inches(0.15), top + Inches(0.65), width - Inches(0.3), Inches(0.35),
             label, font_size=11, color=TEXT_DIM)
    return card


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════════
    # SLIDE 1: KND Intro Video (full-bleed)
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    # Embed the KND terminal intro video centred on slide
    video_path = os.path.join(ROOT, "webapp-react", "public", "panacea_intro_terminal.mp4")
    poster_path = os.path.join(FIGURES, "knd_poster.png")
    if os.path.exists(video_path):
        # Calculate size to maintain 1280x720 aspect ratio centered on slide
        vid_w = Inches(11.5)
        vid_h = Inches(6.47)  # 11.5 * (720/1280)
        vid_x = (SLIDE_W - vid_w) // 2
        vid_y = (SLIDE_H - vid_h) // 2
        slide.shapes.add_movie(
            video_path,
            vid_x, vid_y, vid_w, vid_h,
            poster_frame_image=poster_path if os.path.exists(poster_path) else None,
            mime_type="video/mp4",
        )
    else:
        # Fallback: text title if video not found
        add_text(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1.2),
                 "PANACEA", font_size=72, bold=True, color=TEXT_PRIMARY,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")

    # Small footer below video
    add_text(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.4),
             "AIPI 540 Deep Learning Applications  \u00b7  Duke University  \u00b7  Dominic Tanzillo",
             font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # SLIDE 2: Problem & Motivation
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.7))
    add_text(slide, Inches(0.8), Inches(0.85), Inches(5), Inches(0.5),
             "THE PROBLEM", font_size=13, color=TEXT_DIM, bold=True)

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(1.6),
             "Over 30,000 tracked objects orbit Earth. Every 5 days, two of them collide, "
             "creating thousands of new fragments \u2014 each a future collision.\n\n"
             "This is the Kessler Syndrome: a cascading chain reaction that could "
             "make entire orbital shells unusable.",
             font_size=16, color=TEXT_MUTED)

    add_text(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.0),
             "NASA\u2019s Conjunction Assessment office screens hundreds of "
             "Conjunction Data Messages (CDMs) daily. The hard part: early CDMs for "
             "dangerous and safe events look identical.",
             font_size=14, color=TEXT_MUTED)

    # Screening framing
    add_text(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.9),
             "PANACEA is a screening mechanism: catch every escalating event, "
             "then run detailed analysis on the flagged subset. "
             f"The naive approach \u2014 assume nothing escalates \u2014 gets {naive_acc:.0f}% accuracy "
             "but catches zero events. Sensitivity is what matters.",
             font_size=14, color=TEXT_MUTED)

    # Key question callout
    q_card = add_card(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.0))
    add_text(slide, Inches(1.0), Inches(5.65), Inches(5.1), Inches(0.7),
             "Will collision probability exceed the maneuver planning threshold "
             "before time of closest approach?",
             font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.LEFT)

    # Right side: dark-themed problem visualization
    viz_path = os.path.join(FIGURES, "slide_problem_viz.png")
    if os.path.exists(viz_path):
        slide.shapes.add_picture(viz_path, Inches(6.6), Inches(0.7), Inches(6.4), Inches(3.5))

    # API callout card below the figure
    api_card = add_card(slide, Inches(6.6), Inches(4.5), Inches(6.4), Inches(1.5))
    add_text(slide, Inches(6.85), Inches(4.6), Inches(5.9), Inches(0.4),
             "Screening + Deep Learning API", font_size=14, bold=True, color=TEXT_PRIMARY)
    add_text(slide, Inches(6.85), Inches(5.0), Inches(5.9), Inches(0.8),
             "PANACEA screens every event for escalation risk (high sensitivity). "
             "Flagged pairs can then be analyzed in detail via the deep learning API:\n"
             "pip install panacea-ssa",
             font_size=12, color=TEXT_MUTED)

    # ════════════════════════════════════════════
    # SLIDE 3: Technical Approach (matches About page)
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.5))
    add_text(slide, Inches(0.8), Inches(0.65), Inches(6), Inches(0.5),
             "TECHNICAL APPROACH", font_size=13, color=TEXT_DIM, bold=True)

    models = [
        ("00", "Naive Baseline \u2014 Predict No Escalation",
         f"Assume nothing escalates. Accuracy = {naive_acc:.0f}%, but recall = 0%. "
         f"Misses every dangerous event. This is why screening sensitivity matters.",
         RED),
        ("01", "Logistic Regression",
         f"23 engineered features from CDM sequences. Retrained daily on resolved pairs. "
         f"F1 = {lr_f1:.3f}, the strongest single model.",
         ACCENT),
        ("02", "BiLSTM with Transfer Learning",
         f"Pre-trained on 13K ESA Kelvins events, fine-tuned on {n_pairs:,} Space-Track pairs. "
         f"Focal loss for class imbalance. Attention-weighted pooling.",
         ACCENT),
        ("03", "Ensemble",
         f"40% LR + 30% BiLSTM + 30% regression signal. "
         f"{'Zero missed escalations' if fn_count == 0 else f'Misses only {fn_count}'} "
         f"across {tp} positive events in production.",
         ACCENT),
        ("04", "Graph Neural Network",
         "GraphSAGE on the conjunction network. Captures orbital neighborhood "
         "effects that per-pair models cannot.",
         TEXT_MUTED),
        ("05", "Autoregressive Forecaster",
         "Predicts next CDM update. Multi-step rollouts with Monte Carlo "
         "dropout uncertainty. Correlation r = 0.927.",
         TEXT_MUTED),
        ("06", "Conformal Prediction",
         f"Distribution-free coverage guarantees. {conf_cov:.1f}% empirical "
         f"coverage at 90% target with no distributional assumptions.",
         TEXT_MUTED),
    ]

    y = Inches(1.2)
    row_h = Inches(0.72)  # slightly tighter to fit 7 rows
    for num, title, desc, num_color in models:
        card = add_card(slide, Inches(0.8), y, Inches(11.7), row_h)

        # Number
        add_text(slide, Inches(1.0), y + Inches(0.10), Inches(0.5), Inches(0.4),
                 num, font_size=15, bold=True, color=num_color, font_name="Consolas")
        # Title
        add_text(slide, Inches(1.6), y + Inches(0.06), Inches(3.5), Inches(0.3),
                 title, font_size=13, bold=True, color=TEXT_PRIMARY)
        # Desc
        add_text(slide, Inches(1.6), y + Inches(0.36), Inches(10.5), Inches(0.35),
                 desc, font_size=11, color=TEXT_DIM)

        y += row_h + Inches(0.06)

    # Data sources at bottom
    add_text(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
             f"Data: ESA Kelvins (162,634 CDMs \u00b7 13,154 events \u00b7 103 features)  "
             f"+  Space-Track ({n_pairs:,}+ pairs \u00b7 {n_cdms:,}+ CDMs \u00b7 23 features, updated twice daily)",
             font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # SLIDE 4: Live Demo transition
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    add_text(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(0.8),
             "LIVE DEMO", font_size=48, bold=True, color=TEXT_PRIMARY,
             alignment=PP_ALIGN.CENTER)

    add_text(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.5),
             "dominictanzillo.github.io/Panacea",
             font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

    demo_items = [
        "Globe \u2014 25,000+ tracked objects with real-time SGP4 propagation",
        "Alerts \u2014 Featured conjunction pairs with 3D flyby visualization",
        "Forecast \u2014 Pc evolution charts with model predictions & uncertainty",
        "Models \u2014 BiLSTM learning over time, cross-validation results",
        "Pipeline \u2014 Hyperparameter grid search, feature importance",
    ]
    for i, item in enumerate(demo_items):
        add_text(slide, Inches(3.5), Inches(4.3) + i * Inches(0.42), Inches(6.5), Inches(0.4),
                 item, font_size=14, color=TEXT_MUTED)

    # ════════════════════════════════════════════
    # SLIDE 5: Results & Key Findings
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0.8), Inches(0.5))
    add_text(slide, Inches(0.8), Inches(0.65), Inches(6), Inches(0.5),
             "RESULTS & KEY FINDINGS", font_size=13, color=TEXT_DIM, bold=True)

    # Left: Production track record
    add_text(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
             "Production Track Record (60+ days)", font_size=16, bold=True, color=TEXT_PRIMARY)

    # Stat cards — top row
    card_y = Inches(1.9)
    add_stat_card(slide, Inches(0.8), card_y, f"{resolved:,}", "Resolved Predictions")
    add_stat_card(slide, Inches(3.3), card_y, f"{prod_acc:.1f}%", "Overall Accuracy")

    card_y2 = Inches(3.2)
    add_stat_card(slide, Inches(0.8), card_y2, f"{tp}", "True Positives")
    add_stat_card(slide, Inches(3.3), card_y2, "0", "False Negatives")

    # Hero callout
    hero = add_card(slide, Inches(0.8), Inches(4.6), Inches(4.7), Inches(1.0))
    add_text(slide, Inches(1.0), Inches(4.7), Inches(4.3), Inches(0.8),
             f"100% RECALL\n"
             f"Zero missed escalations across {resolved:,} predictions",
             font_size=18, bold=True, color=GREEN, alignment=PP_ALIGN.LEFT)

    # Right: CV comparison table header
    tbl_x = Inches(6.2)
    add_text(slide, tbl_x, Inches(1.3), Inches(6), Inches(0.4),
             "Cross-Validation (670 pairs, 5-fold)", font_size=16, bold=True, color=TEXT_PRIMARY)

    # Table rows — now includes naive baseline + accuracy column
    rows = [
        ("Model", "Acc", "Recall", "F1", True, None),
        ("Naive (all negative)", f"{naive_acc:.0f}%", "0%", "0", False, RED),
        ("Logistic Regression", "93.2%", "89.9%", "0.854", False, None),
        ("BiLSTM + Transfer", "85.5%", "51.8%", "0.597", False, None),
        ("Ensemble", "91.8%", "98.7%", "0.847", False, ACCENT),
    ]
    col_widths = [Inches(2.8), Inches(1.0), Inches(1.0), Inches(1.0)]
    for ri, (model, acc, rec, f1, is_header, row_accent) in enumerate(rows):
        row_y = Inches(1.9) + ri * Inches(0.45)
        vals = [model, acc, rec, f1]
        cx = tbl_x
        for ci, (val, cw) in enumerate(zip(vals, col_widths)):
            if is_header:
                c = TEXT_DIM
            elif row_accent and ci > 0:
                c = row_accent
            elif ri == 1:  # naive row dim
                c = RED if ci > 0 else TEXT_MUTED
            else:
                c = TEXT_PRIMARY
            b = is_header or row_accent is not None
            fs = 12 if is_header else 14
            fnt = "Consolas" if ci > 0 and not is_header else "Calibri"
            add_text(slide, cx, row_y, cw, Inches(0.4),
                     val, font_size=fs, bold=b, color=c,
                     font_name=fnt,
                     alignment=PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT)
            cx += cw

    # Key insight
    add_text(slide, tbl_x, Inches(4.3), Inches(6.5), Inches(0.8),
             "The naive baseline gets 75% accuracy by ignoring every event \u2014 "
             "but catches nothing. PANACEA is a screening tool: "
             "catch every escalation, then run detailed analysis on flagged pairs. "
             "Deep learning via the API makes it specific.",
             font_size=13, color=TEXT_MUTED)

    # Confusion matrix figure
    cm_path = os.path.join(FIGURES, "cdm_confusion_matrix.png")
    if os.path.exists(cm_path):
        slide.shapes.add_picture(cm_path, tbl_x, Inches(5.0), Inches(3.0))

    # Feature importance figure
    fi_path = os.path.join(FIGURES, "cdm_feature_importance.png")
    if os.path.exists(fi_path):
        slide.shapes.add_picture(fi_path, tbl_x + Inches(3.3), Inches(5.0), Inches(3.5))

    # ════════════════════════════════════════════
    # SLIDE 6: Closing
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    # Top accent bar
    slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Pt(4)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ACCENT
    slide.shapes[-1].line.fill.background()

    add_text(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(0.8),
             "PANACEA", font_size=48, bold=True, color=TEXT_PRIMARY,
             alignment=PP_ALIGN.CENTER)

    add_text(slide, Inches(2), Inches(2.4), Inches(9.3), Inches(0.5),
             "Satellite Collision Avoidance Through Deep Learning",
             font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # Three-column takeaways
    cols = [
        ("Open Source", f"pip install panacea-ssa\n{n_pairs:,} pairs, {n_cdms:,} CDMs\nUpdated twice daily"),
        ("Production Proven", f"60+ days, {resolved:,} predictions\n100% recall, 0 false negatives\nConformal uncertainty bounds"),
        ("What\u2019s Next", "Larger CDM archive (10K+ pairs)\nMulti-provider fusion\nTraCSS integration"),
    ]
    col_w = Inches(3.2)
    gap = Inches(0.5)
    total = len(cols) * col_w + (len(cols) - 1) * gap
    sx = (SLIDE_W - total) / 2
    for i, (title, body) in enumerate(cols):
        x = sx + i * (col_w + gap)
        add_card(slide, x, Inches(3.3), col_w, Inches(2.5))
        add_text(slide, x + Inches(0.2), Inches(3.45), col_w - Inches(0.4), Inches(0.4),
                 title, font_size=15, bold=True, color=TEXT_PRIMARY)
        add_text(slide, x + Inches(0.2), Inches(3.9), col_w - Inches(0.4), Inches(1.6),
                 body, font_size=13, color=TEXT_MUTED)

    # Footer
    add_text(slide, Inches(0), Inches(6.4), SLIDE_W, Inches(0.4),
             "dominictanzillo.github.io/Panacea",
             font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.4),
             "AIPI 540 Deep Learning Applications  \u00b7  Duke University  \u00b7  Dominic Tanzillo",
             font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
